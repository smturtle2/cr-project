"""Fix: rename 모듈 1 → 모듈 3 by handling split runs."""
from pptx import Presentation

PPTX_PATH = "docs/교수님미팅_0410.pptx"
prs = Presentation(PPTX_PATH)

for i in range(3, 9):  # slides 4-9
    slide = prs.slides[i]
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # " 1 " → " 3 " (handles split run case)
                    if run.text == " 1 – Baseline ":
                        run.text = " 3 – Baseline "
                    elif run.text == " 1 – Cross Attention ":
                        run.text = " 3 – Cross Attention "
                    elif run.text == " 1 – ":
                        run.text = " 3 – "
                    elif "모듈1" in run.text:
                        run.text = run.text.replace("모듈1", "모듈3")
                    elif "모듈 1" in run.text:
                        run.text = run.text.replace("모듈 1", "모듈 3")

prs.save(PPTX_PATH)
print("Rename complete.")
