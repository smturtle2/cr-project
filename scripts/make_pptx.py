"""Rewrite module 3 slides: concise problem/solution, one-page architecture, results."""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PPTX_PATH = "docs/교수님미팅_0410.pptx"
prs = Presentation(PPTX_PATH)

BLUE = RGBColor(0x00, 0x70, 0xC0)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x80, 0x00)
ORANGE = RGBColor(0xE0, 0x7C, 0x00)
GRAY = RGBColor(0x99, 0x99, 0x99)

def tb(slide, l, t, w, h, text, sz=14, b=False, c=DARK, al=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = al
    r = p.runs[0]; r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = "맑은 고딕"
    return txBox

def ap(tf, text, sz=14, b=False, c=DARK, al=PP_ALIGN.LEFT, sp=Pt(3)):
    p = tf.add_paragraph(); p.text = text; p.alignment = al
    if sp: p.space_before = sp
    if p.runs:
        r = p.runs[0]; r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = "맑은 고딕"

def box(slide, l, t, w, h, lines, lc=BLUE, fc=WHITE, sz=9, tc=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    s.line.color.rgb = lc; s.line.width = Pt(1.2); s.fill.solid(); s.fill.fore_color.rgb = fc
    tf = s.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(0)
        if p.runs:
            r = p.runs[0]; r.font.size = Pt(sz + 1 if i == 0 else sz); r.font.name = "맑은 고딕"
            r.font.bold = (i == 0); r.font.color.rgb = (tc or lc) if i == 0 else DARK
    return s

def arrow(slide, l, t, text="→", sz=13, c=DARK):
    tb(slide, l, t, 220000, 300000, text, sz=sz, b=True, c=c, al=PP_ALIGN.CENTER)

def hline(slide, l, t, w):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(35000))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()

def vline(slide, l, t, h, c=GRAY):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(25000), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def style_cell(cell, text, sz=11, b=False, c=DARK, bg=None, al=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]; run = p.add_run(); run.text = text; p.alignment = al
    run.font.size = Pt(sz); run.font.bold = b; run.font.color.rgb = c; run.font.name = "맑은 고딕"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        sf = tcPr.makeelement(qn('a:solidFill'), {})
        sf.append(sf.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % bg}))
        tcPr.append(sf)

def clean(slide):
    for s in list(slide.shapes):
        s._element.getparent().remove(s._element)


# ============================================================
# Slide 5: 문제점 + 해결 (간결)
# ============================================================
slide5 = prs.slides[4]
clean(slide5)

tb(slide5, 853440, 426720, 10363200, 600000, "모듈 3 – 문제점 & CAFM", sz=24, b=True, c=BLUE)
hline(slide5, 853440, 1080000, 2194560)

# 문제 (좌)
tb(slide5, 853440, 1250000, 5000000, 350000, "문제  |  SAR 정보 희석 + 균일 처리", sz=15, b=True, c=RED)

bx = slide5.shapes.add_textbox(Emu(853440), Emu(1700000), Emu(5000000), Emu(1600000))
tf = bx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "SAR(2ch) + Opt(13ch) concat → SAR 비중 13%"
r = p.runs[0]; r.font.size = Pt(13); r.font.color.rgb = DARK; r.font.name = "맑은 고딕"
ap(tf, "ResBlock 거치며 SAR 구조 정보 점점 소실", sz=13, c=DARK, sp=Pt(6))
ap(tf, "구름 두꺼운 곳 / 맑은 곳 구분 없이 동일 처리", sz=13, c=DARK, sp=Pt(6))
ap(tf, "", sz=8)
ap(tf, "→ 구름 영역에 적응적 복원 불가", sz=14, b=True, c=RED, sp=Pt(6))

# 구분선
vline(slide5, 6050000, 1250000, 2200000)

# 해결 (우)
tb(slide5, 6350000, 1250000, 5200000, 350000, "해결  |  CAFM", sz=15, b=True, c=GREEN)

bx2 = slide5.shapes.add_textbox(Emu(6350000), Emu(1700000), Emu(5200000), Emu(1600000))
tf2 = bx2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]; p2.text = "A. 구름 밀도 추정 (cosine similarity)"
r2 = p2.runs[0]; r2.font.size = Pt(13); r2.font.bold = True; r2.font.color.rgb = BLUE; r2.font.name = "맑은 고딕"
ap(tf2, "  SAR-Opt 유사도 → density ∈ [0,1]", sz=12, c=DARK, sp=Pt(4))
ap(tf2, "", sz=6)
ap(tf2, "B. 밀도 기반 Feature 변조 (AdaIN)", sz=13, b=True, c=BLUE, sp=Pt(4))
ap(tf2, "  out = feat × (1+γ) + β", sz=12, c=ORANGE, b=True, sp=Pt(4))
ap(tf2, "  맑은→보존 (γ≈0)  /  구름→재구성 (γ>0)", sz=12, c=DARK, sp=Pt(4))
ap(tf2, "", sz=6)
ap(tf2, "Zero-init: 학습 초기 identity → 기존 성능 유지", sz=12, b=True, c=GREEN, sp=Pt(4))

# 하단 핵심 수식 박스
formula = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Emu(853440), Emu(3700000), Emu(10500000), Emu(650000))
formula.fill.solid(); formula.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
formula.line.color.rgb = BLUE; formula.line.width = Pt(1.5)
tf_f = formula.text_frame; tf_f.word_wrap = True
p_f = tf_f.paragraphs[0]; p_f.alignment = PP_ALIGN.CENTER
r_f = p_f.add_run()
r_f.text = "density = Sigmoid(Refine(cos_sim(Enc(SAR), Enc(Opt))))     →     output = feature × (1+γ(density)) + β(density)"
r_f.font.size = Pt(13); r_f.font.bold = True; r_f.font.color.rgb = BLUE; r_f.font.name = "맑은 고딕"


# ============================================================
# Slide 6: 아키텍처 (한 페이지, 간결)
# ============================================================
slide6 = prs.slides[5]
clean(slide6)

tb(slide6, 853440, 350000, 10363200, 550000, "모듈 3 – CAFM 아키텍처", sz=24, b=True, c=BLUE)
hline(slide6, 853440, 950000, 2194560)

# ---- Main pipeline (한 줄) ----
Y = 1300000; H = 580000; gap = 80000

specs = [
    (350000,  850000,  ["Input", "SAR+Opt"],      DARK,  WHITE),
    (1400000, 950000,  ["Head", "Conv 15→256"],    BLUE,  WHITE),
    (2550000, 1050000, ["Body1", "Res×8+Att"],     BLUE,  WHITE),
    (3800000, 850000,  ["CAFM①", "γ₁, β₁"],       ORANGE, RGBColor(0xFD,0xF2,0xE9)),
    (4850000, 1050000, ["Body2", "Res×3+Att"],     BLUE,  WHITE),
    (6100000, 850000,  ["CAFM②", "γ₂, β₂"],       ORANGE, RGBColor(0xFD,0xF2,0xE9)),
    (7150000, 950000,  ["Body3", "Res×4"],          BLUE,  WHITE),
    (8300000, 850000,  ["Tail", "256→13ch"],        BLUE,  WHITE),
]

for x, w, lines, lc, fc in specs:
    box(slide6, x, Y, w, H, lines, lc=lc, fc=fc, sz=9)

# Arrows
arrow_xs = [1200000, 2350000, 3600000, 4650000, 5900000, 6950000, 8100000]
for ax in arrow_xs:
    arrow(slide6, ax, Y + 130000, "→", sz=12)

# Output
tb(slide6, 9150000, Y + 50000, 500000, 250000, "+cloudy", sz=10, b=True, c=DARK, al=PP_ALIGN.CENTER)
tb(slide6, 9150000, Y + 300000, 700000, 300000, "= Output", sz=13, b=True, c=BLUE, al=PP_ALIGN.CENTER)

# ---- Density path (아래) ----
DY = Y + H + 300000

box(slide6, 350000, DY, 700000, 450000, ["SAR 2ch"], lc=ORANGE, fc=RGBColor(0xFF,0xF3,0xE0), sz=9)
arrow(slide6, 1050000, DY + 70000)
box(slide6, 1300000, DY, 800000, 450000, ["Opt 13ch"], lc=BLUE, fc=RGBColor(0xE8,0xF0,0xFE), sz=9)
arrow(slide6, 2100000, DY + 70000)
box(slide6, 2400000, DY, 2200000, 450000, ["CloudDensityEstimator", "cos_sim → refine → σ"], lc=RED, fc=RGBColor(0xFF,0xEB,0xEE), sz=8)
arrow(slide6, 4600000, DY + 70000)
box(slide6, 4900000, DY, 1000000, 450000, ["density", "∈ [0,1]"], lc=RED, fc=RGBColor(0xFF,0xEB,0xEE), sz=9)

# Up arrows to CAFMs
tb(slide6, 4050000, Y + H - 80000, 350000, 400000, "↑", sz=14, b=True, c=RED, al=PP_ALIGN.CENTER)
tb(slide6, 6300000, Y + H - 80000, 350000, 400000, "↑", sz=14, b=True, c=RED, al=PP_ALIGN.CENTER)
tb(slide6, 5900000, DY + 70000, 800000, 300000, "공유", sz=10, b=True, c=RED, al=PP_ALIGN.CENTER)

# ---- CAFM 내부 구조 (하단) ----
CY = DY + 650000
tb(slide6, 853440, CY, 10400000, 350000, "CAFM 내부 흐름", sz=14, b=True, c=ORANGE)

CY2 = CY + 400000
box(slide6, 853440, CY2, 1100000, 500000, ["Feature", "(B,256,H,W)"], lc=BLUE, sz=8)
arrow(slide6, 1953440, CY2 + 100000)
box(slide6, 2200000, CY2, 900000, 500000, ["GAP", "글로벌 벡터"], lc=DARK, sz=8)

# + density
tb(slide6, 3100000, CY2 + 350000, 600000, 250000, "+ density", sz=8, b=True, c=RED, al=PP_ALIGN.CENTER)

arrow(slide6, 3100000, CY2 + 100000)
box(slide6, 3400000, CY2 - 50000, 1000000, 300000, ["Scale Net→γ"], lc=GREEN, fc=RGBColor(0xE8,0xF5,0xE9), sz=9)
box(slide6, 3400000, CY2 + 300000, 1000000, 300000, ["Shift Net→β"], lc=GREEN, fc=RGBColor(0xE8,0xF5,0xE9), sz=9)
arrow(slide6, 4400000, CY2 + 100000)

# Formula
fb = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Emu(4700000), Emu(CY2), Emu(2500000), Emu(500000))
fb.fill.solid(); fb.fill.fore_color.rgb = RGBColor(0xFD,0xF2,0xE9)
fb.line.color.rgb = ORANGE; fb.line.width = Pt(2)
tf_fb = fb.text_frame; tf_fb.word_wrap = True
p_fb = tf_fb.paragraphs[0]; p_fb.alignment = PP_ALIGN.CENTER
r_fb = p_fb.add_run()
r_fb.text = "out = feat×(1+γ) + β"
r_fb.font.size = Pt(14); r_fb.font.bold = True; r_fb.font.color.rgb = ORANGE; r_fb.font.name = "맑은 고딕"

arrow(slide6, 7200000, CY2 + 100000)
box(slide6, 7500000, CY2, 1300000, 500000, ["Modulated", "Feature"], lc=BLUE, sz=9)

# Key point
tb(slide6, 853440, CY2 + 650000, 10400000, 300000,
   "zero-init: γ=0, β=0 → 초기 identity  |  ~100K params (본체 대비 <1%)  |  추가 VRAM 무시 가능",
   sz=10, b=True, c=GRAY, al=PP_ALIGN.CENTER)


# ============================================================
# Slide 7: L1 결과 (기존 slide 9 → slide 7로 이동)
# ============================================================
# slide 7 = prs.slides[6], but currently slides 6(old detail), 7(old arch), 8(old full arch) exist
# We'll write into slide index 6 and delete 7,8
slide7 = prs.slides[6]
clean(slide7)

tb(slide7, 853440, 350000, 10363200, 550000, "모듈 3 – CAFM 실험 결과", sz=24, b=True, c=BLUE)
hline(slide7, 853440, 950000, 2194560)

tb(slide7, 853440, 1050000, 10363200, 300000,
   "L1 Loss  |  4096 train  |  20 epochs  |  batch 4  |  crop 128×128", sz=12, c=DARK)

# Val Best
tb(slide7, 853440, 1450000, 5200000, 300000, "Val Best", sz=14, b=True, c=BLUE)
tbl_v = prs.slides[6].shapes.add_table(3, 5, Emu(853440), Emu(1800000), Emu(5200000), Emu(950000)).table
for c, v in enumerate(("", "MAE↓", "PSNR↑", "SSIM↑", "SAM↓")):
    style_cell(tbl_v.cell(0, c), v, sz=10, b=True, c=WHITE, bg=(0x00,0x70,0xC0))
for c, v in enumerate(("Baseline", "0.2024", "26.37", "0.8418", "9.008")):
    style_cell(tbl_v.cell(1, c), v, sz=10, b=(c==0), c=DARK)
for c, v in enumerate(("CAFM", "0.2027", "26.34", "0.8391", "8.858")):
    style_cell(tbl_v.cell(2, c), v, sz=10, b=(c==0), c=DARK, bg=(0xE8,0xF0,0xFE))

# Test
tb(slide7, 6300000, 1450000, 5000000, 300000, "Test", sz=14, b=True, c=BLUE)
tbl_t = prs.slides[6].shapes.add_table(3, 5, Emu(6300000), Emu(1800000), Emu(5000000), Emu(950000)).table
for c, v in enumerate(("", "MAE↓", "PSNR↑", "SSIM↑", "SAM↓")):
    style_cell(tbl_t.cell(0, c), v, sz=10, b=True, c=WHITE, bg=(0x00,0x70,0xC0))
for c, v in enumerate(("Baseline", "0.1610", "26.93", "0.8584", "9.533")):
    style_cell(tbl_t.cell(1, c), v, sz=10, b=(c==0), c=DARK)
for c, v in enumerate(("CAFM", "0.1591", "26.95", "0.8634", "9.900")):
    style_cell(tbl_t.cell(2, c), v, sz=10, b=(c==0), c=DARK, bg=(0xE8,0xF0,0xFE))

# Improvement
tb(slide7, 853440, 2950000, 10363200, 300000, "CAFM 개선폭 (Test)", sz=13, b=True, c=BLUE)
tbl_i = prs.slides[6].shapes.add_table(2, 4, Emu(853440), Emu(3300000), Emu(10363200), Emu(600000)).table
for c, v in enumerate(("MAE↓", "PSNR↑", "SSIM↑", "SAM↓")):
    style_cell(tbl_i.cell(0, c), v, sz=11, b=True, c=WHITE, bg=(0x00,0x70,0xC0))
vals = [("-0.0019", GREEN), ("+0.02 dB", GREEN), ("+0.005", GREEN), ("+0.367", RED)]
for c, (v, clr) in enumerate(vals):
    style_cell(tbl_i.cell(1, c), v, sz=12, b=True, c=clr, bg=(0xE8,0xF0,0xFE))

# Analysis
tb(slide7, 853440, 4100000, 10363200, 300000, "분석", sz=14, b=True, c=BLUE)

bx_a = slide7.shapes.add_textbox(Emu(853440), Emu(4450000), Emu(10363200), Emu(1400000))
tf_a = bx_a.text_frame; tf_a.word_wrap = True
p_a = tf_a.paragraphs[0]
p_a.text = "MAE/PSNR/SSIM 3개 지표 소폭 개선, SAM만 악화 (+0.37)"
r_a = p_a.runs[0]; r_a.font.size = Pt(12); r_a.font.color.rgb = DARK; r_a.font.name = "맑은 고딕"; r_a.font.bold = True
ap(tf_a, "CAFM best epoch=2 vs Baseline epoch=7 → 수렴 속도 3.5배 빠름", sz=11, c=DARK, sp=Pt(5))
ap(tf_a, "CAFM 단독 효과는 미미 → 추가 모듈 결합 or 대규모 데이터 재검증 필요", sz=11, c=DARK, sp=Pt(5))


# ============================================================
# Delete slides 8, 9 (old architecture and old results)
# ============================================================
for i in range(len(prs.slides) - 1, 6, -1):  # delete from back, indices 8,7
    rId = prs.slides._sldIdLst[i].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[i]
    prs.slides._sldIdLst.remove(sldId)


prs.save(PPTX_PATH)
print("Done! Saved to", PPTX_PATH)
