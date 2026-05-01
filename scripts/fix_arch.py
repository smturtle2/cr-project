"""Rewrite slides 5(problem/solution) and 6(architecture) with IDEA_REPORT content."""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PPTX_PATH = "docs/교수님미팅_0410.pptx"
prs = Presentation(PPTX_PATH)

BLUE = RGBColor(0x00, 0x70, 0xC0)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE0, 0x7C, 0x00)
RED = RGBColor(0xC0, 0x00, 0x00)
GRAY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x00, 0x80, 0x00)
BG_ORANGE = RGBColor(0xFD, 0xF2, 0xE9)
BG_RED = RGBColor(0xFF, 0xEB, 0xEE)
BG_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
BG_GREEN = RGBColor(0xE8, 0xF5, 0xE9)

def tb(slide, l, t, w, h, text, sz=14, b=False, c=DARK, al=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = al
    r = p.runs[0]; r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = "맑은 고딕"
    return bx

def ap(tf, text, sz=11, b=False, c=DARK, al=PP_ALIGN.LEFT, sp=Pt(3)):
    p = tf.add_paragraph(); p.text = text; p.alignment = al; p.space_before = sp
    if p.runs:
        r = p.runs[0]; r.font.size = Pt(sz); r.font.name = "맑은 고딕"
        r.font.bold = b; r.font.color.rgb = c

def box(slide, l, t, w, h, lines, lc=BLUE, fc=WHITE, sz=10, lw=1.5):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    s.line.color.rgb = lc; s.line.width = Pt(lw); s.fill.solid(); s.fill.fore_color.rgb = fc
    tf = s.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(0)
        if p.runs:
            r = p.runs[0]; r.font.size = Pt(sz + 1 if i == 0 else sz)
            r.font.name = "맑은 고딕"; r.font.bold = (i == 0)
            r.font.color.rgb = lc if i == 0 else DARK
    return s

def rect(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def hline(slide, l, t, w, c=BLUE):
    rect(slide, l, t, w, 30000, c)

def vline(slide, l, t, h, c=GRAY):
    rect(slide, l, t, 25000, h, c)

def arrow_h(slide, l, t, sz=13, c=DARK):
    tb(slide, l, t, 200000, 280000, "→", sz=sz, b=True, c=c, al=PP_ALIGN.CENTER)

def arrow_v(slide, l, t, text="↑", sz=14, c=DARK):
    tb(slide, l, t, 250000, 300000, text, sz=sz, b=True, c=c, al=PP_ALIGN.CENTER)

def clean(slide):
    for s in list(slide.shapes):
        s._element.getparent().remove(s._element)


# ============================================================
# Slide 5: 문제 정의 + 해결 아이디어 (IDEA_REPORT 내용 반영)
# ============================================================
slide5 = prs.slides[4]
clean(slide5)

tb(slide5, 700000, 300000, 10500000, 550000,
   "모듈 3 – Cloud-Adaptive Feature Modulation (CAFM)", sz=22, b=True, c=BLUE)
hline(slide5, 700000, 880000, 2500000)

# ---- 문제 정의 (좌측) ----
tb(slide5, 700000, 1050000, 5200000, 350000, "문제 정의", sz=17, b=True, c=RED)

bx_p = slide5.shapes.add_textbox(Emu(700000), Emu(1450000), Emu(5200000), Emu(3200000))
tf_p = bx_p.text_frame; tf_p.word_wrap = True

p = tf_p.paragraphs[0]
p.text = "하나의 위성 이미지 안에서 구름 두께/분포는 매우 불균일"
r = p.runs[0]; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = DARK; r.font.name = "맑은 고딕"

ap(tf_p, "", sz=4)
ap(tf_p, "맑은 영역: 구름 없음, 지표면 정보 온전", sz=12, c=GREEN, b=True)
ap(tf_p, "얇은 구름: 반투명, 지표 일부 가시", sz=12, c=ORANGE, b=True)
ap(tf_p, "두꺼운 구름: 지표 완전 차단, Optical 정보 0%", sz=12, c=RED, b=True)
ap(tf_p, "", sz=4)
ap(tf_p, "→ 각 영역은 서로 다른 복원 전략이 필요", sz=13, b=True, c=RED)
ap(tf_p, "", sz=4)
ap(tf_p, "기존 Baseline (ACA-CRNet)의 한계:", sz=13, b=True, c=DARK)
ap(tf_p, "  SAR(2ch)+Opt(13ch) 단순 concat → SAR 비중 13%", sz=11, c=DARK)
ap(tf_p, "  모든 영역을 동일 가중치로 처리 (one-size-fits-all)", sz=11, c=DARK)
ap(tf_p, "  구름 밀도 정보를 어디에서도 활용하지 않음", sz=11, c=DARK)

# ---- 구분선 ----
vline(slide5, 6100000, 1050000, 3600000)

# ---- 해결 아이디어 (우측) ----
tb(slide5, 6400000, 1050000, 5200000, 350000, "해결 아이디어 — CAFM", sz=17, b=True, c=BLUE)

bx_s = slide5.shapes.add_textbox(Emu(6400000), Emu(1450000), Emu(5200000), Emu(3200000))
tf_s = bx_s.text_frame; tf_s.word_wrap = True

p2 = tf_s.paragraphs[0]
p2.text = "픽셀별 구름 밀도를 추정하고, feature를 적응적 변조"
r2 = p2.runs[0]; r2.font.size = Pt(13); r2.font.bold = True; r2.font.color.rgb = BLUE; r2.font.name = "맑은 고딕"

ap(tf_s, "", sz=4)
ap(tf_s, "1. Cloud Density Estimation", sz=13, b=True, c=BLUE)
ap(tf_s, "  SAR feat · Opt feat 코사인 유사도 기반", sz=11, c=DARK)
ap(tf_s, "  맑은 영역: SAR≈Opt → 유사도 높음 → density≈0", sz=11, c=GREEN)
ap(tf_s, "  구름 영역: SAR≠Opt → 유사도 낮음 → density≈1", sz=11, c=RED)
ap(tf_s, "  근거: SMDCNet(ISPRS'25) — SAR-Opt 유사도 검증", sz=10, c=GRAY)
ap(tf_s, "", sz=4)
ap(tf_s, "2. AdaIN-style Feature Modulation", sz=13, b=True, c=BLUE)
ap(tf_s, "  output = feat × (1 + γ) + β", sz=12, b=True, c=ORANGE)
ap(tf_s, "  맑은 → γ≈0, β≈0 → feature 유지 (원본 보존)", sz=11, c=GREEN)
ap(tf_s, "  얇은 → 중간 변조 (Opt+SAR 블렌딩)", sz=11, c=ORANGE)
ap(tf_s, "  두꺼운 → 강한 변조 (SAR 기반 재구성)", sz=11, c=RED)
ap(tf_s, "  근거: MWFormer(TIP'24), EMRDM(CVPR'25)", sz=10, c=GRAY)

# ---- 하단: 핵심 요약 ----
sum_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Emu(700000), Emu(4900000), Emu(10900000), Emu(600000))
sum_box.fill.solid(); sum_box.fill.fore_color.rgb = BG_BLUE
sum_box.line.color.rgb = BLUE; sum_box.line.width = Pt(1.5)
tf_sum = sum_box.text_frame; tf_sum.word_wrap = True
p_sum = tf_sum.paragraphs[0]; p_sum.alignment = PP_ALIGN.CENTER
r_sum = p_sum.add_run()
r_sum.text = "핵심: 구름 밀도별 적응적 처리 — 맑은 영역은 보존, 구름 영역은 SAR 기반 재구성"
r_sum.font.size = Pt(13); r_sum.font.bold = True; r_sum.font.color.rgb = BLUE; r_sum.font.name = "맑은 고딕"
ap(tf_sum, "zero-init (EMRDM 패턴): γ=0, β=0 → 학습 초기 CAFM = identity, 기존 네트워크 성능 보장",
   sz=11, b=False, c=DARK, al=PP_ALIGN.CENTER, sp=Pt(2))

# ---- 참고 논문 ----
tb(slide5, 700000, 5650000, 10900000, 300000,
   "근거: SMDCNet(ISPRS'25) 유사도 검증 | GLGF-CR(PR'25) 구름 두께별 게이팅 | DFPIR(CVPR'25) 열화 인지 변조 +0.9dB | EMRDM(CVPR'25) zero-init",
   sz=9, c=GRAY)


# ============================================================
# Slide 6: 아키텍처 (IDEA_REPORT 구조 반영)
# ============================================================
slide6 = prs.slides[5]
clean(slide6)

tb(slide6, 600000, 250000, 10500000, 500000,
   "모듈 3 – CAFM 아키텍처", sz=22, b=True, c=BLUE)
hline(slide6, 600000, 780000, 2200000)

# Legend
box(slide6, 600000, 900000, 180000, 180000, [""], lc=BLUE, fc=BG_BLUE, sz=1, lw=1)
tb(slide6, 800000, 900000, 1200000, 200000, "기존 Baseline", sz=10, b=True, c=BLUE)
box(slide6, 2200000, 900000, 180000, 180000, [""], lc=ORANGE, fc=BG_ORANGE, sz=1, lw=1)
tb(slide6, 2400000, 900000, 1200000, 200000, "CAFM (신규)", sz=10, b=True, c=ORANGE)
box(slide6, 3800000, 900000, 180000, 180000, [""], lc=RED, fc=BG_RED, sz=1, lw=1)
tb(slide6, 4000000, 900000, 1500000, 200000, "Density 경로", sz=10, b=True, c=RED)

# ============================================================
# Main pipeline
# ============================================================
Y1 = 1400000; H1 = 700000

box(slide6, 200000, Y1, 850000, H1, ["Input", "Opt 13ch", "SAR 2ch"], lc=DARK, fc=WHITE, sz=9)
arrow_h(slide6, 1050000, Y1 + 200000)

box(slide6, 1300000, Y1, 950000, H1, ["Head", "Concat→Conv", "15→256ch"], lc=BLUE, fc=BG_BLUE, sz=9)
arrow_h(slide6, 2250000, Y1 + 200000)

box(slide6, 2500000, Y1, 1150000, H1, ["Body 1", "ResBlock ×8", "RACAB ×1"], lc=BLUE, fc=BG_BLUE, sz=9)
arrow_h(slide6, 3650000, Y1 + 200000)

# CAFM ①
box(slide6, 3900000, Y1, 1050000, H1, ["CAFM ①", "density →", "γ₁·β₁ 변조"],
    lc=ORANGE, fc=BG_ORANGE, sz=9, lw=2.5)
arrow_h(slide6, 4950000, Y1 + 200000)

box(slide6, 5200000, Y1, 1150000, H1, ["Body 2", "ResBlock ×3", "RACAB ×1"], lc=BLUE, fc=BG_BLUE, sz=9)
arrow_h(slide6, 6350000, Y1 + 200000)

# CAFM ②
box(slide6, 6600000, Y1, 1050000, H1, ["CAFM ②", "density →", "γ₂·β₂ 변조"],
    lc=ORANGE, fc=BG_ORANGE, sz=9, lw=2.5)
arrow_h(slide6, 7650000, Y1 + 200000)

box(slide6, 7900000, Y1, 1050000, H1, ["Body 3", "ResBlock ×4"], lc=BLUE, fc=BG_BLUE, sz=9)
arrow_h(slide6, 8950000, Y1 + 200000)

box(slide6, 9200000, Y1, 850000, H1, ["Tail", "Conv", "256→13"], lc=BLUE, fc=BG_BLUE, sz=9)
arrow_h(slide6, 10050000, Y1 + 100000, sz=11)

tb(slide6, 10300000, Y1 + 30000, 1200000, 250000, "+ cloudy", sz=10, b=True, c=DARK)
tb(slide6, 10300000, Y1 + 300000, 1200000, 350000, "= Output", sz=13, b=True, c=BLUE)

# Global skip
skip_y = Y1 - 180000
rect(slide6, 500000, skip_y, 9900000, 18000, GRAY)
rect(slide6, 500000, skip_y, 18000, 180000, GRAY)
rect(slide6, 10382000, skip_y, 18000, 180000, GRAY)
tb(slide6, 4200000, skip_y - 220000, 2800000, 220000, "Global Residual (cloudy skip)", sz=9, c=GRAY, al=PP_ALIGN.CENTER)

# ============================================================
# Density estimation path
# ============================================================
Y2 = Y1 + H1 + 450000; H2 = 550000

box(slide6, 1300000, Y2, 700000, H2, ["SAR", "2ch"], lc=ORANGE, fc=RGBColor(0xFF, 0xF3, 0xE0), sz=10)
box(slide6, 2100000, Y2, 700000, H2, ["Optical", "13ch"], lc=BLUE, fc=BG_BLUE, sz=10)
arrow_h(slide6, 2800000, Y2 + 120000)

box(slide6, 3050000, Y2, 2500000, H2,
    ["CloudDensityEstimator", "Enc(SAR)·Enc(Opt) cosine sim", "→ Refine → Sigmoid → density"],
    lc=RED, fc=BG_RED, sz=8, lw=2)
arrow_h(slide6, 5550000, Y2 + 120000)

box(slide6, 5800000, Y2, 1000000, H2, ["density", "(B,1,H,W)", "∈ [0,1]"],
    lc=RED, fc=BG_RED, sz=9, lw=2)

# Density → CAFM connections
conn_top = Y1 + H1
conn_bot = Y2 + H2 // 2

# Vertical lines from density to CAFMs
cafm1_x = 4425000
cafm2_x = 7125000
vline(slide6, cafm1_x, conn_top, conn_bot - conn_top, RED)
vline(slide6, cafm2_x, conn_top, conn_bot - conn_top, RED)

# Horizontal connecting line
rect(slide6, cafm1_x, conn_bot, cafm2_x - cafm1_x, 22000, RED)

# Connect to density box
rect(slide6, 6300000, Y2 + 120000, cafm2_x - 6300000, 22000, RED)

# Up arrows
arrow_v(slide6, cafm1_x - 100000, conn_top - 280000, "↑", sz=13, c=RED)
arrow_v(slide6, cafm2_x - 100000, conn_top - 280000, "↑", sz=13, c=RED)

tb(slide6, 6900000, Y2 + 50000, 1500000, 250000, "공유", sz=10, b=True, c=RED)

# ============================================================
# CAFM 내부 상세 (하단 — IDEA_REPORT 구조 반영)
# ============================================================
Y3 = Y2 + H2 + 350000

# Left: Cloud Density Estimation 설명
box(slide6, 600000, Y3, 3400000, 700000,
    ["Cloud Density Estimation",
     "SAR·Opt → 1×1 Conv → cosine sim",
     "→ Refine(Conv→GELU→Conv→σ)",
     "→ density: 0(맑음) ~ 1(두꺼운 구름)"],
    lc=RED, fc=BG_RED, sz=8, lw=1.5)

# Middle: Arrow
arrow_h(slide6, 4000000, Y3 + 200000)

# Middle: Scale/Shift generation
box(slide6, 4250000, Y3, 2200000, 700000,
    ["Scale(γ) / Shift(β) 생성",
     "GAP(feat) + GAP(density)",
     "→ Linear → γ (zero-init)",
     "→ Linear → β (zero-init)"],
    lc=GREEN, fc=BG_GREEN, sz=8, lw=1.5)

# Right: Arrow
arrow_h(slide6, 6450000, Y3 + 200000)

# Right: AdaIN modulation
box(slide6, 6700000, Y3, 2600000, 700000,
    ["AdaIN-style Modulation",
     "out = feat × (1 + γ) + β",
     "맑은(d≈0): 보존 | 얇은: 블렌딩",
     "두꺼운(d≈1): SAR 기반 재구성"],
    lc=ORANGE, fc=BG_ORANGE, sz=8, lw=1.5)

# Spec bar
tb(slide6, 600000, Y3 + 780000, 10800000, 250000,
   "RACAB 직후 2곳 삽입  |  블록별 독립 γ·β 학습  |  ~100K params (본체 <1%)  |  추가 VRAM 무시  |  zero-init → 초기 identity",
   sz=10, b=True, c=GRAY, al=PP_ALIGN.CENTER)

prs.save(PPTX_PATH)
print("Done!")
